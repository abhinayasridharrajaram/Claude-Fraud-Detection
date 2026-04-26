"""Build outputs/Fraud_Detection_Deck.pptx from existing pipeline outputs."""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import os

# ── Palette ────────────────────────────────────────────────────────────────
BLUE    = RGBColor(0x1F, 0x3A, 0x68)   # deep blue accent
GRAY    = RGBColor(0x44, 0x44, 0x44)   # neutral body text
LGRAY   = RGBColor(0x88, 0x88, 0x88)   # source citations
RED     = RGBColor(0xC0, 0x20, 0x20)   # fraud / warning
WHITE   = RGBColor(0xFF, 0xFF, 0xFF)
LBLUE   = RGBColor(0xD0, 0xD8, 0xEA)   # light-blue fill for header rows
YELLOW  = RGBColor(0xFF, 0xF2, 0xCC)   # highlight best cell

BASE    = r"C:\Users\Abhi\claude-bigquery-copilot"
OUT     = os.path.join(BASE, "outputs", "Fraud_Detection_Deck.pptx")

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]   # completely blank

# ── Helper utilities ────────────────────────────────────────────────────────

def add_rect(slide, l, t, w, h, fill=None, line=None, line_w=Pt(0)):
    from pptx.util import Emu
    shape = slide.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
    shape.line.width = line_w
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if line:
        shape.line.color.rgb = line
    else:
        shape.line.fill.background()
    return shape


def add_text(slide, text, l, t, w, h, size=18, bold=False, color=GRAY,
             align=PP_ALIGN.LEFT, wrap=True):
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = wrap
    tf = txb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color
    return txb


def add_bullet_box(slide, items, l, t, w, h, size=17, color=GRAY,
                   bullet_color=None, indent=0.25):
    from pptx.util import Pt, Inches
    txb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txb.word_wrap = True
    tf = txb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(size)
        run.font.color.rgb = color
    return txb


def title_bar(slide, title_text):
    """Full-width blue bar at top with white title."""
    add_rect(slide, 0, 0, 13.33, 1.1, fill=BLUE)
    add_text(slide, title_text, 0.4, 0.12, 12.5, 0.9,
             size=30, bold=True, color=WHITE)


def source_note(slide, note):
    add_text(slide, note, 0.3, 7.1, 12.7, 0.3,
             size=9, color=LGRAY)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s1 = prs.slides.add_slide(BLANK)
add_rect(s1, 0, 0, 13.33, 7.5, fill=BLUE)            # full blue background

# Main title
add_text(s1, "Unsupervised Fraud Detection",
         0.6, 1.4, 12.1, 1.2, size=38, bold=True, color=WHITE,
         align=PP_ALIGN.CENTER)
add_text(s1, "on the Kaggle Credit Card Dataset",
         0.6, 2.5, 12.1, 1.0, size=32, bold=False, color=WHITE,
         align=PP_ALIGN.CENTER)

# Subtitle rule line (white bar)
add_rect(s1, 3.5, 3.7, 6.3, 0.04, fill=WHITE)

add_text(s1, "Anomaly detection + clustering  ·  0 labels used in training",
         0.6, 3.9, 12.1, 0.6, size=18, color=LBLUE, align=PP_ALIGN.CENTER)

add_text(s1, "Run ID: 2026-04-24_1712  ·  Dataset: Kaggle Credit Card Fraud (ULB)",
         0.6, 5.2, 12.1, 0.5, size=14, color=LGRAY, align=PP_ALIGN.CENTER)
add_text(s1, "Date: 2026-04-24",
         0.6, 5.75, 12.1, 0.4, size=13, color=LGRAY, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Problem
# ══════════════════════════════════════════════════════════════════════════════
s2 = prs.slides.add_slide(BLANK)
add_rect(s2, 0, 0, 13.33, 7.5, fill=WHITE)
title_bar(s2, "The Problem")

bullets = [
    "284,807 transactions over 48 hours — only 492 are fraud (0.17%)",
    "Class labels unavailable at inference time; unsupervised scoring required",
    "Naive 'predict all legit' → 99.83% accuracy, 0% fraud caught",
    "Evaluation via precision@top-k and ROC-AUC against held-out labels",
]
add_bullet_box(s2, bullets, 0.5, 1.3, 7.6, 4.5, size=17)

# ── Simple bar chart via shapes ─────────────────────────────────────────────
BAR_L = 8.6
BAR_T = 1.5
BAR_W = 0.9
MAX_H = 3.6

# Legit bar
add_rect(s2, BAR_L, BAR_T + (MAX_H - MAX_H), BAR_W, MAX_H, fill=BLUE)
add_text(s2, "99.83%", BAR_L - 0.05, BAR_T - 0.35, BAR_W + 0.2, 0.4,
         size=13, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
add_text(s2, "Legitimate", BAR_L - 0.05, BAR_T + MAX_H + 0.05,
         BAR_W + 0.2, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Fraud bar (proportional to 0.17 / 99.83)
fraud_h = MAX_H * (0.17 / 99.83)
FRAUD_L = BAR_L + BAR_W + 0.5
add_rect(s2, FRAUD_L, BAR_T + MAX_H - fraud_h, BAR_W, fraud_h, fill=RED)
add_text(s2, "0.17%", FRAUD_L - 0.05, BAR_T + MAX_H - fraud_h - 0.35,
         BAR_W + 0.2, 0.4, size=13, bold=True, color=RED, align=PP_ALIGN.CENTER)
add_text(s2, "Fraud", FRAUD_L - 0.05, BAR_T + MAX_H + 0.05,
         BAR_W + 0.2, 0.4, size=12, color=GRAY, align=PP_ALIGN.CENTER)

# Axis baseline
add_rect(s2, BAR_L - 0.1, BAR_T + MAX_H, 2.3, 0.02, fill=GRAY)

add_text(s2, "Class distribution — 578:1 imbalance",
         8.0, 5.5, 5.0, 0.4, size=11, color=LGRAY, align=PP_ALIGN.CENTER)

source_note(s2, "Source: outputs/eda_summary.md · outputs/runs/2026-04-24_1712/data_quality_report.md")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Approach
# ══════════════════════════════════════════════════════════════════════════════
s3 = prs.slides.add_slide(BLANK)
add_rect(s3, 0, 0, 13.33, 7.5, fill=WHITE)
title_bar(s3, "Approach")

bullets = [
    "Feature engineering: log1p(Amount), hour_of_day, PCA components V1–V28",
    "StandardScaler applied to all 30 features — Class never seen by models",
    "Models: Isolation Forest · One-Class SVM · Local Outlier Factor",
    "Evaluation: precision@top-1% / top-0.1%, recall@top-1%, ROC-AUC vs held-out Class",
]
add_bullet_box(s3, bullets, 0.5, 1.3, 12.3, 3.2, size=17)

# ── Arrow pipeline diagram ────────────────────────────────────────────────
STEPS = [
    ("Data\n284K rows", BLUE),
    ("Features\nlog / hour /\nV1–V28", BLUE),
    ("Models\nIF · LOF\nOCSVM", BLUE),
    ("Score &\nRank", RED),
]
BOX_W = 2.2
BOX_H = 1.5
GAP   = 0.45
START_L = 0.6
DIAG_T  = 4.6

for i, (label, color) in enumerate(STEPS):
    bx = START_L + i * (BOX_W + GAP)
    add_rect(s3, bx, DIAG_T, BOX_W, BOX_H, fill=color)
    add_text(s3, label, bx, DIAG_T + 0.15, BOX_W, BOX_H - 0.1,
             size=14, bold=False, color=WHITE, align=PP_ALIGN.CENTER)
    if i < len(STEPS) - 1:
        arr_l = bx + BOX_W + 0.05
        add_text(s3, "→", arr_l, DIAG_T + 0.45, GAP + 0.1, 0.6,
                 size=22, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

source_note(s3, "Source: ml/00_build_features.py · ml/01_anomaly_detection.py · CLAUDE.md")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Results
# ══════════════════════════════════════════════════════════════════════════════
s4 = prs.slides.add_slide(BLANK)
add_rect(s4, 0, 0, 13.33, 7.5, fill=WHITE)
title_bar(s4, "Results — Anomaly Detection")

# ── Metrics table ─────────────────────────────────────────────────────────
rows_data = [
    ["Model",              "Prec@top-1%", "Prec@top-0.1%", "Recall@top-1%", "ROC-AUC"],
    ["Isolation Forest",   "0.099",       "0.354",          "0.571",          "0.949"],
    ["One-Class SVM",      "0.097",       "0.140",          "0.563",          "0.945"],
    ["Local Outlier Factor","0.007",      "0.000",          "0.039",          "0.513"],
]

from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

COL_W = [2.3, 1.4, 1.55, 1.5, 1.2]
ROW_H = 0.48
TABLE_L = 0.4
TABLE_T = 1.25

for r_idx, row in enumerate(rows_data):
    for c_idx, cell_text in enumerate(row):
        cx = TABLE_L + sum(COL_W[:c_idx])
        cy = TABLE_T + r_idx * ROW_H

        # background
        if r_idx == 0:
            bg = BLUE
        elif r_idx == 1 and c_idx in (2, 4):   # best cells
            bg = YELLOW
        else:
            bg = RGBColor(0xF5, 0xF5, 0xF5) if r_idx % 2 == 0 else WHITE

        add_rect(s4, cx, cy, COL_W[c_idx], ROW_H, fill=bg,
                 line=RGBColor(0xCC, 0xCC, 0xCC), line_w=Pt(0.5))

        txt_color = WHITE if r_idx == 0 else (RED if r_idx == 3 and c_idx > 0 else GRAY)
        bold = r_idx == 0 or (r_idx == 1 and c_idx in (2, 4))
        add_text(s4, cell_text, cx + 0.05, cy + 0.07,
                 COL_W[c_idx] - 0.1, ROW_H - 0.1,
                 size=13, bold=bold, color=txt_color, align=PP_ALIGN.CENTER)

# ★ best-model callout
add_text(s4, "★ Best: Isolation Forest — 206× lift at top-0.1%, ROC-AUC 0.949",
         0.4, TABLE_T + 4 * ROW_H + 0.1, 7.8, 0.5,
         size=14, bold=True, color=BLUE)

# ── PR-curve image ─────────────────────────────────────────────────────────
img_path = os.path.join(BASE, "outputs", "ml", "anomaly_pr_curves.png")
if os.path.exists(img_path):
    s4.shapes.add_picture(img_path,
                          Inches(8.3), Inches(1.2),
                          Inches(4.7), Inches(5.2))

source_note(s4, "Source: outputs/ml/anomaly_metrics.csv · outputs/ml/anomaly_pr_curves.png")


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Takeaways + What's Next
# ══════════════════════════════════════════════════════════════════════════════
s5 = prs.slides.add_slide(BLANK)
add_rect(s5, 0, 0, 13.33, 7.5, fill=WHITE)
title_bar(s5, "Takeaways & What's Next")

# Headline callout box
add_rect(s5, 0.4, 1.25, 12.5, 0.85, fill=LBLUE,
         line=BLUE, line_w=Pt(1.2))
add_text(s5,
         "Isolation Forest captures 64.6% of supervised performance "
         "with zero labels  ·  ROC-AUC gap vs XGBoost: only 0.031 pts",
         0.55, 1.32, 12.2, 0.72,
         size=16, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

# Key findings
add_text(s5, "Key Findings", 0.5, 2.25, 6.0, 0.4,
         size=16, bold=True, color=BLUE)
findings = [
    "Prec@top-0.1% = 35.4%  →  206× lift over random selection",
    "Cluster 0 (7% of txns): fraud rate 0.624% — 3.6× average, no labels used",
    "Hours 2–4 AM: 24% of all fraud on just 5% of volume",
    "LOF fails (AUC 0.51) — density collapse in 30-dim PCA space",
]
add_bullet_box(s5, findings, 0.5, 2.7, 6.2, 3.2, size=15)

# What's next
add_text(s5, "What's Next", 7.1, 2.25, 5.8, 0.4,
         size=16, bold=True, color=BLUE)
nexts = [
    "Supervised upgrade: see outputs/ml_supervised_comparison.md",
    "Production monitoring: drift detection via data_quality_report.md",
    "Threshold tuning for review-queue capacity: threshold_decisions.md",
]
add_bullet_box(s5, nexts, 7.1, 2.7, 5.8, 3.2, size=15)

# vertical divider
add_rect(s5, 6.85, 2.2, 0.03, 3.8, fill=LGRAY)

source_note(s5, "Source: outputs/ml_supervised_comparison.md · outputs/ml/anomaly_metrics.csv · outputs/report.md")


# ══════════════════════════════════════════════════════════════════════════════
# Save
# ══════════════════════════════════════════════════════════════════════════════
os.makedirs(os.path.join(BASE, "outputs"), exist_ok=True)
prs.save(OUT)
print(f"Saved: {OUT}")
